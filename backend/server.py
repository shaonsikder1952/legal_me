from fastapi import FastAPI, APIRouter, UploadFile, File, HTTPException
from fastapi.responses import FileResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional
import uuid
from datetime import datetime, timezone
from emergentintegrations.llm.chat import LlmChat, UserMessage
from PyPDF2 import PdfReader
import io
import re
from pdf_generator import generate_contract_pdf
from docx import Document
import pytesseract
from PIL import Image
from pdf2image import convert_from_bytes
from pptx import Presentation
from openpyxl import load_workbook
from odf import text as odf_text, teletype
from odf.opendocument import load as odf_load

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

app = FastAPI()
api_router = APIRouter(prefix="/api")

# Law Database
LAW_DATABASE = [
    {
        "id": "mietrecht_1",
        "title": "Mietrecht § 535 BGB",
        "description": "Basic rental law - landlord must provide the property in a usable condition",
        "url": "https://www.gesetze-im-internet.de/bgb/__535.html"
    },
    {
        "id": "mietrecht_2",
        "title": "Mietrecht § 556d BGB",
        "description": "Limits on agent fees (Maklergebühren)",
        "url": "https://www.gesetze-im-internet.de/bgb/__556d.html"
    },
    {
        "id": "arbeitsrecht_1",
        "title": "Arbeitsrecht § 611a BGB",
        "description": "Employment contract basics",
        "url": "https://www.gesetze-im-internet.de/bgb/__611a.html"
    },
    {
        "id": "arbeitsrecht_2",
        "title": "Kündigungsschutzgesetz",
        "description": "Protection against dismissal",
        "url": "https://www.gesetze-im-internet.de/kschg/"
    },
    {
        "id": "steuerrecht_1",
        "title": "Einkommensteuergesetz",
        "description": "Income tax law",
        "url": "https://www.gesetze-im-internet.de/estg/"
    },
    {
        "id": "agb_1",
        "title": "AGB-Recht § 307 BGB",
        "description": "Control of standard contract terms - unfair terms are invalid",
        "url": "https://www.gesetze-im-internet.de/bgb/__307.html"
    }
]

# Scam Detection Patterns
SCAM_PATTERNS = [
    {
        "pattern": r"(pay|zahlen|send|überweisen|transfer).{0,100}(advance|voraus|upfront|sofort|immediately|western union|gift card|bitcoin|crypto)",
        "indicator": "Advance payment request",
        "severity": "high"
    },
    {
        "pattern": r"(lottery|gewinn|prize|preis|inheritance|erbe|million|jackpot).{0,100}(won|gewonnen|claim|anspruch)",
        "indicator": "Lottery/prize scam",
        "severity": "high"
    },
    {
        "pattern": r"(urgent|dringend|immediately|sofort|act now|limited time|befristet).{0,100}(action|handeln|respond|antworten|expire|ablaufen)",
        "indicator": "Urgency pressure tactic",
        "severity": "medium"
    },
    {
        "pattern": r"(bank account|bankkonto|credit card|kreditkarte|password|passwort|pin|social security|personal information|ssn)",
        "indicator": "Requests sensitive personal information",
        "severity": "high"
    },
    {
        "pattern": r"(nigerian prince|prince|princess|diplomat|government official|minister).{0,100}(money|geld|transfer|fund)",
        "indicator": "Nigerian prince/419 scam pattern",
        "severity": "high"
    },
    {
        "pattern": r"(work from home|heimarbeit|make money fast|schnell geld|guaranteed income|garantiertes einkommen).{0,100}(no experience|keine erfahrung|easy|einfach)",
        "indicator": "Work-from-home scam",
        "severity": "medium"
    },
    {
        "pattern": r"(IRS|tax authority|finanzamt|legal action|rechtliche schritte|arrest|warrant|haftbefehl).{0,100}(unless|außer|payment|zahlung|immediately|sofort)",
        "indicator": "Government impersonation scam",
        "severity": "high"
    },
    {
        "pattern": r"(click here|klicken sie hier|verify account|konto verifizieren|suspended|gesperrt|update information)",
        "indicator": "Phishing attempt",
        "severity": "high"
    },
    {
        "pattern": r"(refund|rückerstattung|overpayment|überzahlung).{0,100}(send back|zurücksenden|return|zurückgeben|difference|differenz)",
        "indicator": "Overpayment scam",
        "severity": "high"
    },
    {
        "pattern": r"(romance|dating|love|liebe).{0,100}(money|geld|help|hilfe|emergency|notfall|hospital|krankenhaus)",
        "indicator": "Romance scam",
        "severity": "high"
    }
]

# Clause Database with patterns
CLAUSE_DATABASE = [
    {
        "pattern": r"(kündigungsfrist|notice period|termination).{0,50}(1 tag|1 day|sofort|immediate)",
        "risk": "violates",
        "explanation": "Unreasonably short termination period",
        "law_ref": "agb_1"
    },
    {
        "pattern": r"(haftung|liability).{0,50}(ausgeschlossen|excluded|keine)",
        "risk": "violates",
        "explanation": "Blanket liability exclusions are typically invalid",
        "law_ref": "agb_1"
    },
    {
        "pattern": r"(kaution|deposit|security).{0,50}([4-9]|[1-9][0-9]).{0,20}(monat|month)",
        "risk": "violates",
        "explanation": "Deposit exceeds legal maximum of 3 months rent",
        "law_ref": "mietrecht_1"
    },
    {
        "pattern": r"(maklergebühr|agent fee|commission).{0,50}mieter|tenant.{0,50}pay",
        "risk": "attention",
        "explanation": "Tenant may not have to pay agent fees under certain circumstances",
        "law_ref": "mietrecht_2"
    },
    {
        "pattern": r"(probezeit|probation).{0,50}([7-9]|[1-9][0-9]).{0,20}(monat|month)",
        "risk": "violates",
        "explanation": "Probation period exceeds legal maximum of 6 months",
        "law_ref": "arbeitsrecht_1"
    },
    {
        "pattern": r"(vertrag|contract).{0,50}(automatisch|automatically|auto).{0,50}(verlänger|renew|extend)",
        "risk": "attention",
        "explanation": "Auto-renewal clause - ensure you can cancel in time",
        "law_ref": "agb_1"
    },
    {
        "pattern": r"(kündigung|cancellation).{0,50}(schriftlich|written|mail)",
        "risk": "safe",
        "explanation": "Standard cancellation clause requiring written notice",
        "law_ref": "agb_1"
    },
    {
        "pattern": r"(miete|rent).{0,50}(pünktlich|on time|fällig|due)",
        "risk": "safe",
        "explanation": "Standard payment terms",
        "law_ref": "mietrecht_1"
    }
]

# Comprehensive Trusted Resources Database
TRUSTED_LINKS_DATABASE = {
    "rental": {
        "authorities": [
            {"name": "Tenant Protection Association (Mieterschutzbund)", "url": "https://www.mieterschutzbund.de"},
            {"name": "Berlin Tenant Advisory Service", "url": "https://www.berlin.de/sen/stadtentwicklung/wohnen/mieterschutz/"},
            {"name": "Hamburg Tenant Advisory", "url": "https://www.hamburg.de/mieterberatung/"}
        ],
        "alternatives": [
            {"name": "Safer rental listings on ImmobilienScout24", "url": "https://www.immobilienscout24.de"},
            {"name": "Fair rental contract templates", "url": "https://www.mietrecht.de/mustervertrag/"},
            {"name": "Rental law information portal", "url": "https://www.mietrecht.de"}
        ],
        "report": [
            {"name": "Report unfair rental practices", "url": "https://www.verbraucherzentrale.de/beschwerde"},
            {"name": "File complaint with tenant association", "url": "https://www.mieterbund.de/beratung.html"}
        ]
    },
    "employment": {
        "authorities": [
            {"name": "Federal Employment Agency (Bundesagentur für Arbeit)", "url": "https://www.arbeitsagentur.de"},
            {"name": "DGB Labor Union", "url": "https://www.dgb.de"},
            {"name": "Employee Rights Information", "url": "https://www.bmas.de/DE/Arbeit/Arbeitsrecht/arbeitsrecht.html"}
        ],
        "alternatives": [
            {"name": "Fair employment contract templates", "url": "https://www.arbeitsvertrag.org"},
            {"name": "Job search on Federal Employment Agency", "url": "https://www.arbeitsagentur.de/jobsuche/"},
            {"name": "Employment rights guide", "url": "https://www.dgb.de/themen/arbeitsrecht"}
        ],
        "report": [
            {"name": "Report unfair dismissal", "url": "https://www.dgb.de/service/kontakt"},
            {"name": "File workplace complaint", "url": "https://www.bmas.de/DE/Service/Kontakt/kontakt.html"}
        ]
    },
    "immigration": {
        "authorities": [
            {"name": "Immigration Office Berlin", "url": "https://service.berlin.de/dienstleistung/324284/"},
            {"name": "Federal Office for Migration (BAMF)", "url": "https://www.bamf.de"},
            {"name": "Make Your Way in Germany", "url": "https://www.make-it-in-germany.com"}
        ],
        "alternatives": [
            {"name": "Visa and residence permit guide", "url": "https://www.germany.info/us-en/service/visa"},
            {"name": "Integration courses", "url": "https://www.bamf.de/EN/Themen/Integration/integration_node.html"},
            {"name": "Recognition of foreign qualifications", "url": "https://www.anerkennung-in-deutschland.de/html/en/"}
        ],
        "report": [
            {"name": "Immigration consultation services", "url": "https://www.bamf.de/EN/Service/ServiceCenter/servicecenter-node.html"},
            {"name": "Legal advice for migrants", "url": "https://www.diakonie.de/angebote-und-hilfe/migration-und-integration"}
        ]
    },
    "subscription": {
        "authorities": [
            {"name": "Consumer Protection Center (Verbraucherzentrale)", "url": "https://www.verbraucherzentrale.de"},
            {"name": "German Consumer Protection Portal", "url": "https://www.verbraucher.de"}
        ],
        "alternatives": [
            {"name": "Fair subscription comparison", "url": "https://www.vergleich.de"},
            {"name": "Cancel subscriptions properly", "url": "https://www.verbraucherzentrale.de/wissen/vertraege-reklamation/kundenrechte/so-kuendigen-sie-richtig-6892"},
            {"name": "Subscription trap warning list", "url": "https://www.verbraucherzentrale.de/wissen/digitale-welt/onlinedienste/abofallen-im-internet-was-tun-5147"}
        ],
        "report": [
            {"name": "Report subscription scams", "url": "https://www.verbraucherzentrale.de/beschwerde"},
            {"name": "File complaint about unfair contracts", "url": "https://www.bundesnetzagentur.de/DE/Vportal/Verbraucher/Beschwerde/beschwerde-node.html"}
        ]
    },
    "tax": {
        "authorities": [
            {"name": "German Tax Office (Finanzamt)", "url": "https://www.finanzamt.de"},
            {"name": "Federal Ministry of Finance", "url": "https://www.bundesfinanzministerium.de"},
            {"name": "Tax advisor search", "url": "https://www.steuerkanzlei.de"}
        ],
        "alternatives": [
            {"name": "Tax declaration help", "url": "https://www.elster.de"},
            {"name": "Tax calculator", "url": "https://www.bmf-steuerrechner.de"},
            {"name": "Tax deduction guide", "url": "https://www.finanztip.de/steuererklaerung/"}
        ],
        "report": [
            {"name": "Tax consultation services", "url": "https://www.finanzamt.de/beratung"},
            {"name": "Report tax fraud", "url": "https://www.bundesfinanzministerium.de/Web/DE/Service/Kontakt/kontakt.html"}
        ]
    },
    "general": {
        "authorities": [
            {"name": "Consumer Protection Center", "url": "https://www.verbraucherzentrale.de"},
            {"name": "German Legal Portal", "url": "https://www.gesetze-im-internet.de"},
            {"name": "Federal Ministry of Justice", "url": "https://www.bmjv.de"}
        ],
        "alternatives": [
            {"name": "Legal advice directory", "url": "https://www.anwaltauskunft.de"},
            {"name": "Free legal help", "url": "https://www.rechtshilfe.de"},
            {"name": "Official forms and templates", "url": "https://www.formulare-bfinv.de"}
        ],
        "report": [
            {"name": "General complaint portal", "url": "https://www.verbraucherzentrale.de/beschwerde"},
            {"name": "File consumer complaint", "url": "https://www.bundesnetzagentur.de/DE/Vportal/Verbraucher/Beschwerde/beschwerde-node.html"}
        ]
    }
}

# Pydantic Models
class ChatMessage(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    session_id: str
    user_message: str
    ai_response: str
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class ChatRequest(BaseModel):
    session_id: str
    message: str

class ChatResponse(BaseModel):
    response: str
    session_id: str

class ContractAnalysis(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    filename: str
    extracted_text: str
    document_type: str
    risk_level: str
    page_count: int
    clauses_safe: List[dict]
    clauses_attention: List[dict]
    clauses_violates: List[dict]
    summary: str
    recommendations: str
    key_excerpts: List[str]
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

@api_router.get("/")
async def root():
    return {"message": "LegalMe API"}

@api_router.get("/laws")
async def get_laws():
    return LAW_DATABASE

@api_router.get("/topics")
async def get_topics():
    return [
        {"id": "rental", "name": "Rental Law", "icon": "home"},
        {"id": "employment", "name": "Employment Law", "icon": "briefcase"},
        {"id": "subscription", "name": "Subscription Law", "icon": "file-text"},
        {"id": "tax", "name": "Tax Law", "icon": "calculator"}
    ]

@api_router.post("/chat", response_model=ChatResponse)
async def chat(request: ChatRequest):
    try:
        # Load conversation history for this session
        conversation_history = await db.chat_messages.find(
            {"session_id": request.session_id},
            {"_id": 0}
        ).sort("timestamp", 1).to_list(100)
        
        # Build conversation context
        conversation_context = ""
        if conversation_history:
            conversation_context = "\n\nCONVERSATION HISTORY (for context):\n"
            for msg in conversation_history[-10:]:  # Last 10 messages for context
                conversation_context += f"User: {msg['user_message']}\n"
                conversation_context += f"Assistant: {msg['ai_response'][:200]}...\n\n"
        
        # Create system message with law database context and trusted links
        law_context = "\n".join([f"- {law['title']}: {law['description']} (Link: {law['url']})" for law in LAW_DATABASE])
        
        # Create trusted links context for AI
        links_context = """
TRUSTED LINKS BY CATEGORY (Always use these for Next Steps):

RENTAL:
- Authorities: Mieterschutzbund (https://www.mieterschutzbund.de), Berlin Tenant Advisory (https://www.berlin.de/sen/stadtentwicklung/wohnen/mieterschutz/)
- Alternatives: ImmobilienScout24 (https://www.immobilienscout24.de), Rental templates (https://www.mietrecht.de/mustervertrag/)
- Report: Verbraucherzentrale (https://www.verbraucherzentrale.de/beschwerde)

EMPLOYMENT:
- Authorities: Bundesagentur für Arbeit (https://www.arbeitsagentur.de), DGB Union (https://www.dgb.de)
- Alternatives: Employment contract templates (https://www.arbeitsvertrag.org), Job search (https://www.arbeitsagentur.de/jobsuche/)
- Report: DGB contact (https://www.dgb.de/service/kontakt)

IMMIGRATION:
- Authorities: Immigration Berlin (https://service.berlin.de/dienstleistung/324284/), BAMF (https://www.bamf.de)
- Alternatives: Visa guide (https://www.germany.info/us-en/service/visa), Integration courses (https://www.bamf.de/EN/Themen/Integration/)
- Report: BAMF Service (https://www.bamf.de/EN/Service/ServiceCenter/servicecenter-node.html)

SUBSCRIPTION:
- Authorities: Verbraucherzentrale (https://www.verbraucherzentrale.de)
- Alternatives: Cancel guide (https://www.verbraucherzentrale.de/wissen/vertraege-reklamation/kundenrechte/so-kuendigen-sie-richtig-6892)
- Report: Consumer complaint (https://www.verbraucherzentrale.de/beschwerde)

TAX:
- Authorities: Finanzamt (https://www.finanzamt.de), Tax advisor (https://www.steuerkanzlei.de)
- Alternatives: ELSTER tax portal (https://www.elster.de), Tax calculator (https://www.bmf-steuerrechner.de)
- Report: Tax consultation (https://www.finanzamt.de/beratung)

GENERAL FALLBACK:
- Verbraucherzentrale (https://www.verbraucherzentrale.de)
- Legal portal (https://www.gesetze-im-internet.de)
- Legal advice (https://www.anwaltauskunft.de)
"""
        
        system_message = f"""You are LegalMe, a professional German legal assistant. 

CRITICAL: LINK FORMATTING RULES (NO EXCEPTIONS):
- ALWAYS use markdown link format: [Blue Clickable Text](URL)
- NEVER show raw URLs
- NEVER use HTML <a> tags
- Every law reference MUST be a masked link
- Example: [§ 535 BGB – Rental Agreements](https://www.gesetze-im-internet.de/bgb/__535.html)

FORMATTING RULES:
- Use # ## ### for headers
- Use **bold** for key points
- Use bullet lists with -
- Use --- for horizontal dividers

OFFICIAL GERMAN LAW SOURCES (ONLY USE THESE):
1. Gesetze im Internet: https://www.gesetze-im-internet.de/
   - BGB (Civil Code): https://www.gesetze-im-internet.de/bgb/
   - StGB (Criminal Code): https://www.gesetze-im-internet.de/stgb/
   - AufenthG (Residence Act): https://www.gesetze-im-internet.de/aufenthg_2004/
   - EStG (Income Tax): https://www.gesetze-im-internet.de/estg/
   - Format: https://www.gesetze-im-internet.de/[LAW_CODE]/__[SECTION].html

2. Official German Laws Available:
{law_context}

HOW TO DYNAMICALLY GENERATE LAW LINKS:
- Rental issues → § 535-580 BGB → https://www.gesetze-im-internet.de/bgb/__[section].html
- Employment issues → § 611a-630 BGB → https://www.gesetze-im-internet.de/bgb/__[section].html
- Consumer rights → § 312-312m BGB → https://www.gesetze-im-internet.de/bgb/__[section].html
- Criminal law → § 1-358 StGB → https://www.gesetze-im-internet.de/stgb/__[section].html
- Immigration → § 1-104 AufenthG → https://www.gesetze-im-internet.de/aufenthg_2004/__[section].html

KEY BGB SECTIONS YOU MUST KNOW:
- § 535: Basic rental obligations
- § 551: Rental deposit (max 3 months)
- § 558: Rent increases
- § 573: Termination by landlord
- § 611a: Employment contract basics
- § 622: Notice periods for employment
- § 626: Extraordinary termination
- § 312g: Right of withdrawal
- § 307: Unfair contract terms

{links_context}

MANDATORY REQUIREMENTS FOR EVERY RESPONSE:

1. IDENTIFY RELEVANT GERMAN LAWS:
   - Analyze the user's legal question
   - Identify specific BGB/StGB/AufenthG sections
   - Generate the official gesetze-im-internet.de URL
   - Format as masked markdown link: [§ XXX Law Name](URL)

2. INCLUDE "Relevant Laws" SECTION:
   Every response MUST end with this section (before Next Steps):

   ---
   
   ## Relevant Laws
   - [§ XXX BGB – Description](https://www.gesetze-im-internet.de/bgb/__XXX.html)
   - [§ YYY BGB – Description](https://www.gesetze-im-internet.de/bgb/__YYY.html)

3. IF NO EXACT LAW EXISTS:
   State: "No specific law covers this exactly, but the closest framework is..."
   Then link to the general legal area

4. DETECT USER INTENT and provide context-aware Next Steps:
   - Rental issues → Mieterschutzbund, rental templates
   - Employment → Arbeitsagentur, DGB
   - Immigration → BAMF, Immigration office
   - Subscription → Verbraucherzentrale, cancellation guides

5. END YOUR RESPONSE WITH:

---

## Relevant Laws
- [§ XXX LAW – Description](official_url)
- [§ YYY LAW – Description](official_url)

---

## Next Steps
### 1. [Context-aware action]
[Masked Link Name](URL)

### 2. [Context-aware resource]
[Masked Link Name](URL)

### 3. Upload another document
[Analyze another contract](/contract)

EXAMPLE for rental deposit question:
User: "My landlord wants 4 months deposit"

Response:
This is illegal in Germany. According to § 551 BGB, rental deposits are limited to a maximum of 3 months' cold rent...

---

## Relevant Laws
- [§ 551 BGB – Rental Deposit Limits](https://www.gesetze-im-internet.de/bgb/__551.html)
- [§ 307 BGB – Unfair Contract Terms](https://www.gesetze-im-internet.de/bgb/__307.html)

---

## Next Steps
### 1. Get legal help
[Tenant Protection Association](https://www.mieterschutzbund.de)

### 2. Report this violation
[Consumer Protection Center](https://www.verbraucherzentrale.de/beschwerde)

### 3. Upload another document
[Analyze another contract](/contract)

Be professional, concise, and ALWAYS include masked law links.

CONVERSATION MEMORY:
You are continuing a conversation with this user. Here is the recent history:
{conversation_context}

CRITICAL MEMORY RULES:
- This is session ID: {request.session_id}
- ONLY use conversation history from THIS session
- DO NOT mix up different chat sessions
- Remember what the user asked before in THIS conversation
- Reference previous topics from THIS session if relevant
- Don't repeat information already provided in THIS session
- Maintain continuity based on THIS conversation's context
- If this is a new session with no history, treat it as a fresh conversation"""
        
        # Initialize chat with memory
        chat_client = LlmChat(
            api_key=os.environ['EMERGENT_LLM_KEY'],
            session_id=request.session_id,
            system_message=system_message
        )
        chat_client.with_model("gemini", "gemini-2.0-flash")
        
        # Send message
        user_msg = UserMessage(text=request.message)
        ai_response = await chat_client.send_message(user_msg)
        
        # Store in database
        chat_doc = ChatMessage(
            session_id=request.session_id,
            user_message=request.message,
            ai_response=ai_response
        )
        doc = chat_doc.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        await db.chat_messages.insert_one(doc)
        
        return ChatResponse(response=ai_response, session_id=request.session_id)
    except Exception as e:
        logging.error(f"Chat error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

def chunk_text(text: str, chunk_size: int = 3000) -> list:
    """Split text into manageable chunks for analysis"""
    words = text.split()
    chunks = []
    current_chunk = []
    current_length = 0
    
    for word in words:
        current_length += len(word) + 1
        if current_length > chunk_size:
            chunks.append(' '.join(current_chunk))
            current_chunk = [word]
            current_length = len(word)
        else:
            current_chunk.append(word)
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks

def extract_text_from_file(content: bytes, filename: str) -> tuple:
    """
    Extract text from ANY file type (PDF, DOCX, TXT, Images, XLSX, PPTX, etc.)
    Returns: (extracted_text, page_count)
    Handles files of any size with chunking and OCR
    """
    file_ext = filename.lower().split('.')[-1]
    
    try:
        # PDF files
        if file_ext == 'pdf':
            pdf_reader = PdfReader(io.BytesIO(content))
            page_count = len(pdf_reader.pages)
            extracted_text = ""
            
            # Try regular text extraction first
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    extracted_text += page_text + "\n"
            
            # If no text extracted, it's likely a scanned PDF - use OCR
            if len(extracted_text.strip()) < 50:
                logging.info("PDF has no text, attempting OCR...")
                images = convert_from_bytes(content)
                extracted_text = ""
                for i, image in enumerate(images):
                    text = pytesseract.image_to_string(image)
                    extracted_text += f"\n--- Page {i+1} ---\n{text}\n"
                logging.info(f"OCR extracted {len(extracted_text)} characters")
            
            return extracted_text, page_count
        
        # DOCX files
        elif file_ext in ['docx', 'doc']:
            doc = Document(io.BytesIO(content))
            extracted_text = "\n".join([para.text for para in doc.paragraphs])
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        extracted_text += "\n" + cell.text
            return extracted_text, len(doc.paragraphs) // 20 or 1
        
        # Excel files
        elif file_ext in ['xlsx', 'xls']:
            wb = load_workbook(io.BytesIO(content), data_only=True)
            extracted_text = ""
            for sheet in wb.worksheets:
                extracted_text += f"\n--- Sheet: {sheet.title} ---\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell else "" for cell in row])
                    extracted_text += row_text + "\n"
            return extracted_text, len(wb.worksheets)
        
        # PowerPoint files
        elif file_ext in ['pptx', 'ppt']:
            prs = Presentation(io.BytesIO(content))
            extracted_text = ""
            for i, slide in enumerate(prs.slides):
                extracted_text += f"\n--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
            return extracted_text, len(prs.slides)
        
        # OpenDocument files
        elif file_ext in ['odt', 'ods']:
            doc = odf_load(io.BytesIO(content))
            extracted_text = ""
            for para in doc.getElementsByType(odf_text.P):
                extracted_text += teletype.extractText(para) + "\n"
            return extracted_text, 1
        
        # TXT and other text files
        elif file_ext in ['txt', 'log', 'md', 'rtf', 'csv']:
            try:
                extracted_text = content.decode('utf-8', errors='ignore')
            except:
                extracted_text = content.decode('latin-1', errors='ignore')
            return extracted_text, len(extracted_text) // 3000 or 1
        
        # All image files with OCR
        elif file_ext in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'tif', 'gif', 'webp', 'heic', 'heif']:
            logging.info(f"Processing image file with OCR: {file_ext}")
            image = Image.open(io.BytesIO(content))
            # Convert to RGB if needed
            if image.mode in ('RGBA', 'LA', 'P'):
                image = image.convert('RGB')
            extracted_text = pytesseract.image_to_string(image)
            logging.info(f"OCR extracted {len(extracted_text)} characters from image")
            return extracted_text, 1
        
        # Fallback: Try to extract as text or use OCR
        else:
            logging.warning(f"Unknown file type: {file_ext}, attempting extraction...")
            # Try as text first
            try:
                extracted_text = content.decode('utf-8', errors='ignore')
                if len(extracted_text.strip()) > 50:
                    return extracted_text, 1
            except:
                pass
            
            # Try as image with OCR
            try:
                image = Image.open(io.BytesIO(content))
                extracted_text = pytesseract.image_to_string(image)
                if len(extracted_text.strip()) > 10:
                    logging.info(f"Fallback OCR succeeded for {file_ext}")
                    return extracted_text, 1
            except:
                pass
            
            raise HTTPException(status_code=400, detail=f"Could not extract text from {file_ext} file. The file may be corrupted or in an unsupported format.")
    
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Text extraction error for {file_ext}: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Could not extract text from file: {str(e)}")

@api_router.post("/contract/analyze")
async def analyze_contract(file: UploadFile = File(...)):
    try:
        # Extract text from any supported file type
        content = await file.read()
        extracted_text, page_count = extract_text_from_file(content, file.filename)
        
        if not extracted_text.strip():
            raise HTTPException(status_code=400, detail="Could not extract text from file. The file may be empty or corrupted.")
        
        logging.info(f"Extracted {len(extracted_text)} characters from {page_count} pages/sections")
        
        # SCAM DETECTION - Check for scam patterns first
        scam_indicators = []
        text_lower = extracted_text.lower()
        
        for scam_pattern in SCAM_PATTERNS:
            matches = re.finditer(scam_pattern["pattern"], text_lower, re.IGNORECASE)
            for match in matches:
                snippet = extracted_text[max(0, match.start()-100):min(len(extracted_text), match.end()+100)]
                scam_indicators.append({
                    "indicator": scam_pattern["indicator"],
                    "severity": scam_pattern["severity"],
                    "snippet": snippet.strip()
                })
        
        is_likely_scam = len([s for s in scam_indicators if s["severity"] == "high"]) >= 2 or len(scam_indicators) >= 4
        
        if is_likely_scam:
            logging.warning(f"SCAM DETECTED: {len(scam_indicators)} indicators found")
        
        # Split into chunks for analysis if document is large
        text_chunks = chunk_text(extracted_text, chunk_size=3000)
        logging.info(f"Split document into {len(text_chunks)} chunks")
        
        # Analyze clauses across all chunks
        clauses_safe = []
        clauses_attention = []
        clauses_violates = []
        
        for chunk_idx, chunk in enumerate(text_chunks):
            text_lower = chunk.lower()
            for clause_pattern in CLAUSE_DATABASE:
                matches = re.finditer(clause_pattern["pattern"], text_lower, re.IGNORECASE)
                for match in matches:
                    snippet = chunk[max(0, match.start()-50):min(len(chunk), match.end()+50)]
                    law_ref = next((law for law in LAW_DATABASE if law["id"] == clause_pattern.get("law_ref")), None)
                    
                    # Check for duplicates
                    clause_text = snippet.strip()
                    if any(clause_text in c["clause"] for c in clauses_safe + clauses_attention + clauses_violates):
                        continue
                    
                    clause_info = {
                        "clause": clause_text,
                        "explanation": clause_pattern["explanation"],
                        "law": law_ref["title"] if law_ref else "General Legal Principle",
                        "law_link": law_ref["url"] if law_ref else "#"
                    }
                    
                    if clause_pattern["risk"] == "safe":
                        clauses_safe.append(clause_info)
                    elif clause_pattern["risk"] == "attention":
                        clauses_attention.append(clause_info)
                    elif clause_pattern["risk"] == "violates":
                        clauses_violates.append(clause_info)
        
        # Determine overall risk
        if clauses_violates:
            risk_level = "high"
        elif clauses_attention:
            risk_level = "medium"
        else:
            risk_level = "low"
        
        # Generate comprehensive AI analysis using chunks
        law_context = "\n".join([f"- {law['title']}: {law['description']}" for law in LAW_DATABASE])
        
        # Analyze document in chunks and merge results
        chunk_analyses = []
        for i, chunk in enumerate(text_chunks[:5]):  # Analyze first 5 chunks max
            chunk_prompt = f"""Analyze this section of a legal document:

Section {i+1}:
{chunk}

Identify:
1. Document type (rental/employment/subscription/immigration/tax/other)
2. Key terms and conditions
3. Potential risks or concerns
4. Important deadlines or fees mentioned
5. Missing information

Provide brief analysis (2-3 sentences)."""

            chat_client = LlmChat(
                api_key=os.environ['EMERGENT_LLM_KEY'],
                session_id=f"analysis_chunk_{uuid.uuid4()}",
                system_message="You are a legal document analyzer. Be concise and identify key points."
            )
            chat_client.with_model("gemini", "gemini-2.0-flash")
            
            chunk_analysis = await chat_client.send_message(UserMessage(text=chunk_prompt))
            chunk_analyses.append(f"Section {i+1}: {chunk_analysis}")
        
        # Merge all chunk analyses into final summary
        merged_prompt = f"""You analyzed a legal document in {len(chunk_analyses)} sections. Here are the findings:

{chr(10).join(chunk_analyses)}

Detected clauses:
- Safe clauses: {len(clauses_safe)}
- Attention needed: {len(clauses_attention)}
- Violations: {len(clauses_violates)}

Available German laws:
{law_context}

Now provide a comprehensive final analysis in this EXACT format:
TYPE: [rental/employment/subscription/immigration/tax/other]
SUMMARY: [3-5 sentence comprehensive summary covering key points, risks, and overall assessment]
RECOMMENDATIONS: [Specific actionable recommendations based on the entire document]
KEY_EXCERPTS: [3-5 most important text excerpts from the document, each 50-100 words]"""
        
        chat_client = LlmChat(
            api_key=os.environ['EMERGENT_LLM_KEY'],
            session_id=f"contract_{uuid.uuid4()}",
            system_message="You are a professional German legal document analyzer. Provide comprehensive analysis."
        )
        chat_client.with_model("gemini", "gemini-2.0-flash")
        
        ai_analysis = await chat_client.send_message(UserMessage(text=merged_prompt))
        
        # Parse AI response
        doc_type = "general"
        summary = "Document analysis complete."
        recommendations = "Review all highlighted clauses carefully."
        key_excerpts = []
        
        if "TYPE:" in ai_analysis:
            doc_type = ai_analysis.split("TYPE:")[1].split("\n")[0].strip().lower()
        if "SUMMARY:" in ai_analysis:
            summary_text = ai_analysis.split("SUMMARY:")[1]
            if "RECOMMENDATIONS:" in summary_text:
                summary = summary_text.split("RECOMMENDATIONS:")[0].strip()
            else:
                summary = summary_text.split("KEY_EXCERPTS:")[0].strip() if "KEY_EXCERPTS:" in summary_text else summary_text.strip()
        if "RECOMMENDATIONS:" in ai_analysis:
            rec_text = ai_analysis.split("RECOMMENDATIONS:")[1]
            recommendations = rec_text.split("KEY_EXCERPTS:")[0].strip() if "KEY_EXCERPTS:" in rec_text else rec_text.strip()
        if "KEY_EXCERPTS:" in ai_analysis:
            excerpts_text = ai_analysis.split("KEY_EXCERPTS:")[1].strip()
            # Extract key excerpts (limited to keep report concise)
            key_excerpts = [e.strip() for e in excerpts_text.split("\n") if e.strip()][:5]
        
        # Create analysis document
        analysis = ContractAnalysis(
            filename=file.filename,
            extracted_text=extracted_text,
            document_type=doc_type,
            risk_level=risk_level,
            page_count=page_count,
            clauses_safe=clauses_safe,
            clauses_attention=clauses_attention,
            clauses_violates=clauses_violates,
            summary=summary,
            recommendations=recommendations,
            key_excerpts=key_excerpts
        )
        
        # Store in database
        doc = analysis.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        await db.contract_analyses.insert_one(doc)
        
        return analysis
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Contract analysis error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/contract/{contract_id}")
async def get_contract_analysis(contract_id: str):
    analysis = await db.contract_analyses.find_one({"id": contract_id}, {"_id": 0})
    if not analysis:
        raise HTTPException(status_code=404, detail="Contract analysis not found")
    return analysis

@api_router.get("/contract/{contract_id}/download")
async def download_contract_pdf(contract_id: str):
    analysis = await db.contract_analyses.find_one({"id": contract_id}, {"_id": 0})
    if not analysis:
        raise HTTPException(status_code=404, detail="Contract analysis not found")
    
    try:
        pdf_buffer = generate_contract_pdf(analysis)
        
        from fastapi.responses import StreamingResponse
        return StreamingResponse(
            pdf_buffer,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f"attachment; filename=LegalMe_Report_{contract_id[:8]}.pdf"
            }
        )
    except Exception as e:
        logging.error(f"PDF generation error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"PDF generation failed: {str(e)}")

@api_router.get("/alternatives/{category}")
async def get_alternatives(category: str):
    alt = next((a for a in ALTERNATIVE_DATABASE if a["category"] == category), None)
    if not alt:
        return ALTERNATIVE_DATABASE[2]  # Return general resources
    return alt

@api_router.get("/chat/history")
async def get_chat_history():
    """Get all unique chat sessions"""
    try:
        pipeline = [
            {"$sort": {"timestamp": -1}},
            {"$group": {
                "_id": "$session_id",
                "last_message": {"$first": "$user_message"},
                "timestamp": {"$first": "$timestamp"}
            }},
            {"$sort": {"timestamp": -1}},
            {"$limit": 50}
        ]
        
        sessions = await db.chat_messages.aggregate(pipeline).to_list(50)
        
        result = []
        for session in sessions:
            result.append({
                "session_id": session["_id"],
                "preview": session["last_message"][:60] + "..." if len(session["last_message"]) > 60 else session["last_message"],
                "timestamp": session["timestamp"]
            })
        
        return result
    except Exception as e:
        logging.error(f"Chat history error: {str(e)}")
        return []

@api_router.get("/chat/{session_id}/messages")
async def get_session_messages(session_id: str):
    """Get all messages for a specific session"""
    try:
        messages = await db.chat_messages.find(
            {"session_id": session_id},
            {"_id": 0}
        ).sort("timestamp", 1).to_list(1000)
        
        return messages
    except Exception as e:
        logging.error(f"Session messages error: {str(e)}")
        return []

@api_router.put("/chat/{session_id}/rename")
async def rename_session(session_id: str, request: dict):
    """Rename a chat session"""
    try:
        new_name = request.get("name", "")
        if not new_name:
            raise HTTPException(status_code=400, detail="Name is required")
        
        result = await db.chat_messages.update_many(
            {"session_id": session_id},
            {"$set": {"session_name": new_name}}
        )
        
        return {"success": True, "updated": result.modified_count}
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Rename session error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.delete("/chat/{session_id}")
async def delete_session(session_id: str):
    """Delete a chat session"""
    try:
        result = await db.chat_messages.delete_many({"session_id": session_id})
        return {"success": True, "deleted": result.deleted_count}
    except Exception as e:
        logging.error(f"Delete session error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/contract/{contract_id}/chat", response_model=ChatResponse)
async def contract_chat(contract_id: str, request: ChatRequest):
    """Chat about a specific contract"""
    try:
        # Get contract analysis
        analysis = await db.contract_analyses.find_one({"id": contract_id}, {"_id": 0})
        if not analysis:
            raise HTTPException(status_code=404, detail="Contract not found")
        
        # Load conversation history for this contract chat
        chat_history = await db.contract_chats.find(
            {"contract_id": contract_id, "session_id": request.session_id},
            {"_id": 0}
        ).sort("timestamp", 1).to_list(50)
        
        # Build conversation context
        chat_context = ""
        if chat_history:
            chat_context = "\n\nPREVIOUS QUESTIONS ABOUT THIS CONTRACT:\n"
            for msg in chat_history[-5:]:  # Last 5 messages
                chat_context += f"User asked: {msg['user_message']}\n"
                chat_context += f"You answered: {msg['ai_response'][:150]}...\n\n"
        
        # Create context with contract details
        contract_context = f"""
CONTRACT CONTEXT:
Type: {analysis.get('document_type', 'unknown')}
Risk Level: {analysis.get('risk_level', 'unknown')}
Summary: {analysis.get('summary', '')}

Safe Clauses: {len(analysis.get('clauses_safe', []))}
Attention Clauses: {len(analysis.get('clauses_attention', []))}
Violating Clauses: {len(analysis.get('clauses_violates', []))}

Recommendations: {analysis.get('recommendations', '')}

Contract Text (excerpt):
{analysis.get('extracted_text', '')[:1500]}
"""
        
        law_context = "\\n".join([f"- {law['title']}: {law['description']}" for law in LAW_DATABASE])
        
        system_message = f"""You are LegalMe, a professional German legal assistant analyzing a specific contract.

CONTRACT YOU ARE ANALYZING:
{contract_context}

CRITICAL: LINK FORMATTING (MANDATORY):
- ALWAYS use markdown links: [Text](URL)
- NEVER use HTML <a> tags or raw URLs
- Every law MUST be a masked link
- Example: [§ 551 BGB – Rental Deposit](https://www.gesetze-im-internet.de/bgb/__551.html)

OFFICIAL GERMAN LAW SOURCES:
- BGB sections: https://www.gesetze-im-internet.de/bgb/__[section].html
- Available laws: {law_context}

MANDATORY FOR EVERY RESPONSE:
1. Answer the user's question about THIS specific contract
2. Reference exact clauses from the analysis
3. Cite relevant German laws with official links
4. Include "Relevant Laws" section with masked links
5. Professional formatting with headers, bullets

RESPONSE FORMAT:
[Your answer referencing specific contract clauses]

---

## Relevant Laws
- [§ XXX BGB – Description](official_url)

---

## Next Steps
### 1. [Relevant action]
[Masked Link](url)

CONVERSATION MEMORY FOR THIS CONTRACT:
{chat_context}

IMPORTANT:
- Remember what the user already asked about THIS contract
- Don't repeat information from previous answers
- Reference previous Q&A if relevant
- Keep the conversation flowing naturally

User's current question: {request.message}"""
        
        # Initialize chat with memory
        chat_client = LlmChat(
            api_key=os.environ['EMERGENT_LLM_KEY'],
            session_id=f"contract_{contract_id}_{request.session_id}",
            system_message=system_message
        )
        chat_client.with_model("gemini", "gemini-2.0-flash")
        
        # Send message
        user_msg = UserMessage(text=request.message)
        ai_response = await chat_client.send_message(user_msg)
        
        # Store in contract chat history
        chat_doc = {
            "id": str(uuid.uuid4()),
            "contract_id": contract_id,
            "session_id": request.session_id,
            "user_message": request.message,
            "ai_response": ai_response,
            "timestamp": datetime.now(timezone.utc).isoformat()
        }
        await db.contract_chats.insert_one(chat_doc)
        
        return ChatResponse(response=ai_response, session_id=request.session_id)
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Contract chat error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/contract/{contract_id}/chat/history")
async def get_contract_chat_history(contract_id: str, session_id: str):
    """Get chat history for a contract"""
    try:
        messages = await db.contract_chats.find(
            {"contract_id": contract_id, "session_id": session_id},
            {"_id": 0}
        ).sort("timestamp", 1).to_list(1000)
        
        return messages
    except Exception as e:
        logging.error(f"Contract chat history error: {str(e)}")
        return []

app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()